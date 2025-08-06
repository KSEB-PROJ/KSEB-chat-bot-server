"""
LangChain 도구: 동적 PowerPoint(.pptx) 프레젠테이션 생성 (v7: 최종 버전)
"""

import os
import json
import uuid
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt

from langchain.tools import tool
from langchain_core.prompts import ChatPromptTemplate
from langchain_openai import ChatOpenAI

from src.core.config import settings
from src.utils.api_client import fetch_messages_from_backend
from .web_search_tool import DeepSearchTool
from .semantic_scholar_tool import SemanticScholarTool

# --- 초기 설정 ---
llm = ChatOpenAI(
    model=settings.OPENAI_MODEL, temperature=0.4, api_key=settings.OPENAI_API_KEY
)
web_search_tool = DeepSearchTool()
paper_search_tool = SemanticScholarTool()
TEMP_DIR = tempfile.gettempdir()


# --- AI PPT 컨텍스트 생성 로직 (v7: '비평 및 개선' 프레임워크) ---
def create_ppt_context(topic: str, combined_info: str) -> dict:
    """LLM을 사용하여 '비평 및 개선' 가이드라인이 포함된 PPT 구조를 JSON으로 생성합니다."""
    system_prompt = r"""
    당신은 TED 강연자를 코칭하는 세계 최고의 프레젠테이션 코치이자, McKinsey의 수석 컨설턴트입니다.
    당신의 임무는 주어진 주제와 데이터를 바탕으로, 청중을 완벽하게 설득하고 감동시킬 PPT 초안을 'JSON' 형식으로 생성하는 것입니다.

    **프레젠테이션 생성 원칙:**
    1.  **표지 정보:** `presentation_title`, `presenter`, `team_members`, `major` 필드를 생성합니다.
    2.  **동적 목차 생성:** 먼저 본문 슬라이드들을 모두 구상한 뒤, 각 섹션의 제목(`Layout_Section_Header`)을 추출하여 `Layout_TOC` 슬라이드의 `content`를 동적으로 구성해야 합니다.
    3.  **(핵심!) '비평 및 개선' 가이드라인:** 모든 슬라이드의 `guideline`은 아래 3가지 항목을 반드시 포함하는 **'PPT 업그레이드 작업 지시서'**여야 합니다.
        - `[Critique]`: (AI가 방금 만든 슬라이드 초안의 부족한 점, 논리적 허점, 설득력이 약한 부분을 날카롭게 진단하고 비평합니다.)
        - `[Action Items]`: (비평에 대한 구체적인 개선 계획을 2개 이상 제시합니다. '데이터 추가', '사례 보강', '시각 자료 제작' 등 즉시 실행 가능한 작업이어야 합니다.)
        - `[Presentation Tip]`: (개선된 슬라이드를 발표할 때, 청중의 마음을 사로잡을 수 있는 강력한 오프닝 멘트나 발표 전략을 제시합니다.)
    4.  **(중요!) 슬라이드 레이아웃 선택:** 기본적으로 모든 내용은 **`"layout": "Layout_Body"`**를 사용하세요. 하지만, 내용상 다이어그램, 그래프, 차트 등 '시각 자료가 반드시 필요한 경우에만' **`"layout": "Layout_Image_Right"`**를 사용하고, `"image_description"` 필드에 필요한 이미지를 구체적으로 묘사해야 합니다. 일반적인 텍스트 슬라이드에 이미지 레이아웃을 남용하지 마세요.

    **JSON 출력 예시:**
    {{
      "presentation_title": "MCP 서버 기반 초개인화 여행 경험 제안",
      "presenter": "김연아",
      "team_members": "박태환, 손흥민",
      "major": "경영정보시스템",
      "slides": [
        {{
          "layout": "Layout_TOC",
          "title": "Index",
          "content": "- I. The Burning Platform: 현 여행 시장의 위기\n- II. The Silver Bullet: MCP 서버, 유일한 해결책\n- III. The Roadmap: 3단계 실행 계획",
          "guideline": "[Critique]: 목차만으로는 발표의 전체적인 스토리와 기대감이 전달되지 않습니다.\n[Action Items]: 1. (디자인) 각 목차 옆에 연관된 아이콘(돋보기, 열쇠, 지도 등)을 배치하여 시각적 이해를 돕고, 2. (부제 추가) '우리는 어떻게 고객의 마음을 얻을 것인가?' 와 같은 질문 형태의 부제를 추가하여 흥미를 유발하세요.\n[Presentation Tip]: '오늘, 저는 여러분께 단순한 기술이 아닌, 여행의 패러다임을 바꿀 새로운 미래에 대해 말씀드리고자 합니다.' 라는 멘트로 발표를 시작하십시오."
        }},
        {{
          "layout": "Layout_Body",
          "title": "고객들은 이미 정답을 알고 있습니다",
          "content": "- '나만을 위한 여행'에 대한 갈망 폭발\n- 현실은 획일적인 패키지 상품의 반복",
          "guideline": "[Critique]: 주장은 있으나, 이를 뒷받침할 객관적인 데이터가 없어 공허하게 들립니다.\n[Action Items]: 1. (데이터 보강) 구글 트렌드나 관련 논문에서 '개인화 여행' 검색량 추이 데이터를 찾아, '지난 3년간 250% 급증'과 같은 구체적인 수치를 담은 꺾은선 그래프를 추가하세요. 2. (인용 추가) '결국 모든 비즈니스는 개인화로 귀결될 것이다' 와 같은 유명 경영자의 명언을 하단에 작게 추가하여 주장에 무게를 더하세요.\n[Presentation Tip]: '데이터가 보여주듯, 이것은 더 이상 선택이 아닌 생존의 문제입니다.' 와 같이 위기감을 고조시키며 슬라이드를 시작하십시오."
        }},
        {{
          "layout": "Layout_Image_Right",
          "title": "MCP 서버 아키텍처",
          "content": "- 사용자 데이터 분석을 통한 실시간 개인화\n- MSA 기반의 유연한 확장성 확보",
          "image_description": "중앙에 MCP 서버가 위치하고, 좌측에는 데이터 수집 모듈(앱, 웹), 우측에는 개인화된 여행 상품을 제공하는 파트너사 API가 연결된 아키텍처 다이어그램. 데이터 흐름을 화살표로 명확히 표시.",
          "guideline": "[Critique]: 아키텍처가 너무 기술 중심으로만 설명되어 비즈니스 측면의 장점이 드러나지 않습니다.\n[Action Items]: 1. (Benefit 추가) 각 구성요소 옆에 '비용 30% 절감', '전환율 20% 상승' 등 기대효과를 작은 텍스트로 추가하세요. 2. (시각화 강화) 다이어그램에 애니메이션 효과를 적용하여 데이터 흐름을 순차적으로 보여주면 이해도가 높아집니다.\n[Presentation Tip]: '이 복잡해 보이는 그림이, 어떻게 우리에게 수백억의 가치를 안겨줄 수 있는지 지금부터 설명드리겠습니다.' 와 같이 호기심을 자극하는 멘트로 시작하세요."
        }}
      ]
    }}
    """
    human_prompt = "발표 주제: {topic}\n\n[참고 자료]:\n{information}"
    prompt = ChatPromptTemplate.from_messages(
        [("system", system_prompt), ("human", human_prompt)]
    )
    chain = prompt | llm

    try:
        response = chain.invoke({"topic": topic, "information": combined_info})
        json_string = (
            response.content.strip().replace("```json", "").replace("```", "").strip()
        )
        return json.loads(json_string)
    except json.JSONDecodeError:
        print("   - 2.3a: AI가 유효하지 않은 JSON을 생성. 자가 수정을 시도합니다.")
        correction_prompt_text = (
            "The following text is a malformed JSON string. "
            "Please correct the syntax errors and return only the valid JSON object. "
            "Do not add any explanations or surrounding text. Just the corrected JSON.\n\n"
            "Malformed JSON:\n---\n{malformed_json}\n---"
        )
        correction_prompt = ChatPromptTemplate.from_template(correction_prompt_text)
        correction_chain = correction_prompt | llm
        corrected_response = correction_chain.invoke(
            {"malformed_json": response.content}
        )
        try:
            corrected_json_string = (
                corrected_response.content.strip()
                .replace("```json", "")
                .replace("```", "")
                .strip()
            )
            print("   - 2.3b: JSON 자가 수정 성공.")
            return json.loads(corrected_json_string)
        except json.JSONDecodeError as e:
            print(
                f"Fatal: JSON 자가 수정 실패. \n오류: {e}\n수정된 내용: {corrected_response.content}"
            )
            raise ValueError(
                "AI가 PPT 구조를 생성하고 수정하는 데 모두 실패했습니다."
            ) from e


# --- 최종 PPTX 파일 생성 로직 (v7) ---
def build_pptx_from_context(context: dict) -> str:
    """사용자가 정의한 커스텀 레이아웃을 찾아 PPT 문서를 생성하고, 맨 앞에 안내 슬라이드를 추가합니다."""
    template_dir = os.path.join(
        os.path.dirname(__file__), "..", "..", "templates", "pptx"
    )
    template_path = os.path.join(template_dir, "default_template.pptx")

    if not os.path.exists(template_path):
        os.makedirs(template_dir, exist_ok=True)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        try:
            prs.slide_layouts[0].name = "Layout_Title"
            prs.slide_layouts[1].name = "Layout_Body"
            prs.slide_layouts[2].name = "Layout_Section_Header"
            prs.slide_layouts[3].name = "Layout_TOC"
            prs.slide_layouts[5].name = "Layout_Conclusion"
            prs.slide_layouts[6].name = "Layout_QNA"
            prs.slide_layouts[8].name = "Layout_Image_Right"
        except IndexError:
            pass
        prs.save(template_path)

    prs = Presentation(template_path)
    layout_map = {layout.name: layout for layout in prs.slide_layouts}

    # 안내 슬라이드
    guide_layout = layout_map.get("Layout_Body", prs.slide_layouts[1])
    guide_slide = prs.slides.add_slide(guide_layout)
    if guide_slide.shapes.title:
        guide_slide.shapes.title.text = "⚠️ 중요: PPT 디자인 수정 방법 (슬라이드 마스터)"
    body_ph = next(
        (ph for ph in guide_slide.placeholders if ph.placeholder_format.idx > 0), None
    )
    if body_ph:
        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True
        p1 = tf.add_paragraph()
        p1.text = "이 PPT의 배경, 로고, 글꼴 등 공통 디자인 요소는 '슬라이드 마스터'에서 관리됩니다."
        p1.font.size = Pt(18)
        p2 = tf.add_paragraph()
        p2.text = "배경 이미지나 로고 등을 수정/삭제하시려면, 이 파일의 슬라이드 마스터에서 직접 편집해야 합니다."
        p2.font.size = Pt(16)
        p2.level = 1
        p3 = tf.add_paragraph()
        p3.text = "수정 방법: [보기] 탭 → [슬라이드 마스터]"
        p3.font.bold = True
        p3.font.size = Pt(16)
        p3.level = 2
        p4 = tf.add_paragraph()
        p4.text = "\n이 안내 슬라이드는 최종 발표 전에 삭제하시면 됩니다."
        p4.font.italic = True
        p4.font.size = Pt(14)

    # 표지 슬라이드
    title_layout = layout_map.get("Layout_Title", prs.slide_layouts[0])
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = context.get("presentation_title", "제목 없음")
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = (
            f"발표자: {context.get('presenter', '')}\n팀원: {context.get('team_members', '')}\n전공: {context.get('major', '')}"
        )

    # 본문 슬라이드
    for s_data in context.get("slides", []):
        layout_name = s_data.get("layout", "Layout_Body")
        slide_layout = layout_map.get(
            layout_name, layout_map.get("Layout_Body", prs.slide_layouts[1])
        )
        slide = prs.slides.add_slide(slide_layout)

        if slide.shapes.title:
            slide.shapes.title.text = s_data.get("title", "")

        # 레이아웃 유형에 따라 콘텐츠 채우기
        if layout_name == "Layout_Image_Right":
            # 텍스트 플레이스홀더 (보통 idx=1)
            text_ph = next(
                (ph for ph in slide.placeholders if ph.placeholder_format.idx == 1),
                None,
            )
            if text_ph:
                tf = text_ph.text_frame
                tf.clear()
                tf.word_wrap = True
                content = s_data.get("content", "")
                content_lines = (
                    content if isinstance(content, list) else str(content).split("\n")
                )
                for line in content_lines:
                    p = tf.add_paragraph()
                    p.text = str(line).strip("- ").strip()
                    p.level = 1 if str(line).strip().startswith("-") else 0

            # 이미지 플레이스홀더 (보통 idx=2)
            img_ph = next(
                (ph for ph in slide.placeholders if ph.placeholder_format.idx == 2),
                None,
            )
            if img_ph:
                # 실제 이미지를 넣는 대신, 사용자에게 안내 텍스트를 추가
                tf = img_ph.text_frame
                tf.clear()
                p = tf.add_paragraph()
                p.text = "⬇️\n여기에 이미지 추가"
                p.font.bold = True
                p.font.size = Pt(18)
                p = tf.add_paragraph()
                p.text = f"({s_data.get('image_description', '이미지 설명 없음')})"
                p.font.size = Pt(14)
        else:
            # 기존 텍스트 전용 슬라이드 로직 (Layout_Body, Layout_TOC 등)
            body_ph = next(
                (ph for ph in slide.placeholders if ph.placeholder_format.idx > 0), None
            )
            if body_ph:
                tf = body_ph.text_frame
                tf.clear()
                tf.word_wrap = True
                content = s_data.get("content", "")
                content_lines = (
                    content if isinstance(content, list) else str(content).split("\n")
                )
                for i, line in enumerate(content_lines):
                    p = (
                        tf.paragraphs[i]
                        if i < len(tf.paragraphs)
                        else tf.add_paragraph()
                    )
                    p.text = str(line).strip("- ").strip()
                    p.level = 1 if str(line).strip().startswith("-") else 0

        guideline = s_data.get("guideline", "가이드라인 없음")
        if isinstance(guideline, list):
            guideline = "\n".join(guideline)
        slide.notes_slide.notes_text_frame.text = str(guideline)

    # 마지막 Q&A 슬라이드
    qna_layout = layout_map.get("Layout_QNA", prs.slide_layouts[6])
    slide = prs.slides.add_slide(qna_layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Q & A"

    output_filename = f"{uuid.uuid4()}.pptx"
    output_path = os.path.join(TEMP_DIR, output_filename)
    prs.save(output_path)
    return output_filename


@tool
async def generate_ppt(
    topic: str, user_id: int, channel_id: int, jwt_token: str
) -> str:
    """(v7) AI가 PPT 초안과 '비평 및 개선' 가이드라인을 제안하고, 오류 시 자가 수정을 시도하며, 커스텀 레이아웃을 찾아 .pptx 파일을 생성합니다."""
    print("---" + "DYNAMIC PPT GENERATION (Final Version) START" + "---")
    print(f"User: {user_id}, Topic: '{topic}'")

    # 1. 정보 수집
    try:
        chat_history = await fetch_messages_from_backend(channel_id, user_id, jwt_token)
        web_research_results = web_search_tool.run(f"{topic}에 대한 발표 자료")
        paper_results = paper_search_tool.run(f"topic: {topic}, limit=3")
        combined_info = (
            f"발표 주제: {topic}\n\n"
            f"[채널 대화 내용 요약]:\n{chat_history}\n\n"
            f"[웹 리서치 결과]:\n{web_research_results}\n\n"
            f"[학술 논문 정보]:\n{paper_results}"
        )
    except Exception as e:
        print(f"Error during data collection: {e}")
        return "발표 자료를 만들기 위한 정보를 수집하는 데 실패했습니다."

    # 2. AI로 PPT 구조 생성 (자가 수정 기능 포함)
    try:
        ppt_context = create_ppt_context(topic, combined_info)
    except (ValueError, json.JSONDecodeError) as e:
        print(f"Error creating PPT context: {e}")
        return "AI가 PPT 구조를 생성하고 수정하는 데 모두 실패했습니다. 다시 시도해 주세요."

    # 3. PPTX 파일 생성
    try:
        output_filename = build_pptx_from_context(ppt_context)
        print(f"   - 3.2: ✅ PPT 생성 성공: {output_filename}")
    except Exception as e:
        print(f"Fatal: PPTX 파일 생성 중 오류 발생: {e}")
        return "PPTX 파일 생성에 실패했습니다. 템플릿 파일의 슬라이드 마스터를 확인하거나 관리자에게 문의하세요."

    # 4. 다운로드 링크 반환
    download_url = (
        f"{settings.CHATBOT_SERVER_URL}{settings.API_V1_STR}/download/{output_filename}"
    )
    ppt_title = ppt_context.get("presentation_title", topic)
    print("---" + "DYNAMIC PPT GENERATION END" + "---")
    return f"'{ppt_title}' 발표 자료 초안이 완성되었습니다.\n[여기에서 다운로드]({download_url})하여 내용을 확인하고 발전시키세요."
